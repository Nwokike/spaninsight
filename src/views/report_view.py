"""Report view — displays pinned charts + descriptions with export and sharing."""

from __future__ import annotations

import base64
import io
import logging

import flet as ft

from core import theme, tokens
from core.state import state
from components.chart_card import build_chart_card
from services.ad_service import AdService

logger = logging.getLogger(__name__)


def build_report_view(page: ft.Page, on_back=None, ad_service: AdService | None = None) -> ft.View:
    """Build the report view with export and sharing."""

    # ── Export Handlers ──────────────────────────────────────────────

    async def on_export_pdf(e):
        if not state.charts:
            return

        page.snack_bar = ft.SnackBar(ft.Text("Generating PDF..."), duration=2000)
        page.snack_bar.open = True
        page.update()

        try:
            # Show interstitial ad during generation
            if ad_service:
                await ad_service.show_interstitial()

            pdf_bytes = _generate_pdf()

            # Flet 0.85.0: FilePicker is a Service, NOT a Control
            picker = ft.FilePicker()
            result = await picker.save_file(
                dialog_title="Save Report PDF",
                file_name="spaninsight_report.pdf",
                allowed_extensions=["pdf"],
            )
            if result:
                with open(result, "wb") as f:
                    f.write(pdf_bytes)
                page.snack_bar = ft.SnackBar(ft.Text("PDF saved!"), duration=3000)
                page.snack_bar.open = True
                page.update()
        except Exception as err:
            logger.exception("PDF export failed")
            page.snack_bar = ft.SnackBar(
                ft.Text(f"PDF export failed: {err}", color=ft.Colors.WHITE),
                bgcolor=theme.ERROR, duration=4000)
            page.snack_bar.open = True
            page.update()

    async def on_export_pptx(e):
        if not state.charts:
            return

        page.snack_bar = ft.SnackBar(ft.Text("Generating PowerPoint..."), duration=2000)
        page.snack_bar.open = True
        page.update()

        try:
            if ad_service:
                await ad_service.show_interstitial()

            pptx_bytes = _generate_pptx()

            picker = ft.FilePicker()
            result = await picker.save_file(
                dialog_title="Save Report PPTX",
                file_name="spaninsight_report.pptx",
                allowed_extensions=["pptx"],
            )
            if result:
                with open(result, "wb") as f:
                    f.write(pptx_bytes)
                page.snack_bar = ft.SnackBar(ft.Text("PPTX saved!"), duration=3000)
                page.snack_bar.open = True
                page.update()
        except Exception as err:
            logger.exception("PPTX export failed")
            page.snack_bar = ft.SnackBar(
                ft.Text(f"PPTX failed: {err}", color=ft.Colors.WHITE),
                bgcolor=theme.ERROR, duration=4000)
            page.snack_bar.open = True
            page.update()

    async def on_share_report(e):
        if not state.charts:
            return

        page.snack_bar = ft.SnackBar(ft.Text("Uploading report..."), duration=2000)
        page.snack_bar.open = True
        page.update()

        try:
            report_json = _build_report_json()

            import httpx
            from core.constants import API_BASE_URL, APP_SECRET, USER_AGENT

            async with httpx.AsyncClient() as client:
                resp = await client.post(
                    f"{API_BASE_URL}/reports",
                    headers={
                        "X-App-Secret": APP_SECRET,
                        "User-Agent": USER_AGENT,
                        "Content-Type": "application/json",
                    },
                    json={
                        "user_uuid": state.user_uuid,
                        "report_json": report_json,
                    },
                    timeout=15.0,
                )
                if resp.status_code == 201:
                    data = resp.json()
                    url = data.get("url", "")
                    page.clipboard = url
                    page.snack_bar = ft.SnackBar(
                        ft.Text(f"Link copied! {url}"), duration=5000)
                    page.snack_bar.open = True
                    page.update()
                else:
                    page.snack_bar = ft.SnackBar(
                        ft.Text("Upload failed. Try again."), duration=3000)
                    page.snack_bar.open = True
                    page.update()

        except Exception as err:
            logger.exception("Share report failed")
            page.snack_bar = ft.SnackBar(
                ft.Text(f"Share failed: {err}", color=ft.Colors.WHITE),
                bgcolor=theme.ERROR, duration=4000)
            page.snack_bar.open = True
            page.update()

    # ── Export Generators ────────────────────────────────────────────

    def _figure_to_png_bytes(figure) -> bytes:
        """Convert a matplotlib figure to PNG bytes."""
        buf = io.BytesIO()
        figure.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        buf.seek(0)
        return buf.read()

    def _generate_pdf() -> bytes:
        """Generate a PDF report using fpdf2."""
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Title page
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 24)
        pdf.cell(0, 40, "Spaninsight Report", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.set_font("Helvetica", "", 12)
        pdf.cell(0, 10, f"{len(state.charts)} insights generated", new_x="LMARGIN", new_y="NEXT", align="C")

        for i, chart in enumerate(state.charts):
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 14)
            pdf.cell(0, 10, f"{i + 1}. {chart.get('prompt', 'Analysis')[:80]}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(5)

            # Chart image
            if chart.get("figure"):
                try:
                    img_bytes = _figure_to_png_bytes(chart["figure"])
                    img_buf = io.BytesIO(img_bytes)
                    pdf.image(img_buf, x=15, w=180)
                    pdf.ln(5)
                except Exception as err:
                    logger.warning("Could not embed chart %d: %s", i, err)

            # Description
            desc = chart.get("description", chart.get("insight", ""))
            if desc:
                pdf.set_font("Helvetica", "", 11)
                pdf.multi_cell(0, 6, desc)

        # Footer
        pdf.add_page()
        pdf.set_font("Helvetica", "I", 10)
        pdf.cell(0, 10, "Generated by Spaninsight — Privacy-First Data Intelligence", align="C")

        return pdf.output()

    def _generate_pptx() -> bytes:
        """Generate a PowerPoint report using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Spaninsight Report"
        slide.placeholders[1].text = f"{len(state.charts)} insights"

        for i, chart in enumerate(state.charts):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

            # Title
            from pptx.util import Emu
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            tf.text = chart.get("prompt", f"Analysis {i + 1}")
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True

            # Chart
            if chart.get("figure"):
                try:
                    img_bytes = _figure_to_png_bytes(chart["figure"])
                    img_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(img_stream, Inches(1), Inches(1.3), width=Inches(8))
                except Exception:
                    pass

            # Description
            desc = chart.get("description", chart.get("insight", ""))
            if desc:
                txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
                tf2 = txBox2.text_frame
                tf2.text = desc
                tf2.paragraphs[0].font.size = Pt(14)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _build_report_json() -> dict:
        """Build JSON for R2 upload (Base64 chart images + descriptions)."""
        items = []
        for chart in state.charts:
            item = {
                "prompt": chart.get("prompt", ""),
                "description": chart.get("description", chart.get("insight", "")),
            }
            if chart.get("figure"):
                try:
                    img_bytes = _figure_to_png_bytes(chart["figure"])
                    item["image_b64"] = base64.b64encode(img_bytes).decode("utf-8")
                except Exception:
                    pass
            items.append(item)

        return {
            "title": f"Spaninsight Report — {state.current_df_name or 'Analysis'}",
            "chart_count": len(items),
            "items": items,
        }

    # ── Content ─────────────────────────────────────────────────────

    if not state.charts:
        content = ft.Container(
            content=ft.Column([
                ft.Container(height=100),
                ft.Icon(ft.Icons.BAR_CHART_ROUNDED, size=80, color=ft.Colors.with_opacity(0.1, ft.Colors.ON_SURFACE)),
                ft.Text("No report yet", size=16, weight="w500", color=ft.Colors.ON_SURFACE_VARIANT),
                ft.Text("Pin analysis results to build your report.", size=13, color=ft.Colors.with_opacity(0.5, ft.Colors.ON_SURFACE)),
                ft.Container(height=20),
                ft.Button("Start Analysis", icon=ft.Icons.ANALYTICS_ROUNDED,
                          on_click=lambda _: setattr(page, 'route', '/analysis') or page.update()),
            ], horizontal_alignment="center", spacing=8),
            expand=True, alignment=ft.Alignment.CENTER,
        )
    else:
        report_cards = [
            build_chart_card(
                index=i + 1,
                prompt=c.get("prompt", ""),
                figure=c.get("figure"),
                insight=c.get("description", c.get("insight", "")),
                code="",
            )
            for i, c in enumerate(state.charts)
        ]

        content = ft.Column([
            ft.Container(
                content=ft.Text(f"Your report contains {len(state.charts)} insights", size=13, color=ft.Colors.ON_SURFACE_VARIANT),
                padding=ft.Padding(20, 10, 0, 0),
            ),
            ft.Column(controls=report_cards, spacing=16, padding=20),
            ft.Container(height=80),
        ], scroll="auto", expand=True)

    appbar = ft.AppBar(
        title=ft.Text("Report", weight="bold"),
        bgcolor=ft.Colors.TRANSPARENT,
        actions=[
            ft.IconButton(ft.Icons.PICTURE_AS_PDF_ROUNDED, tooltip="Export PDF",
                          on_click=lambda e: page.run_task(on_export_pdf, e),
                          visible=len(state.charts) > 0),
            ft.IconButton(ft.Icons.SLIDESHOW_ROUNDED, tooltip="Export PPTX",
                          on_click=lambda e: page.run_task(on_export_pptx, e),
                          visible=len(state.charts) > 0),
            ft.IconButton(ft.Icons.SHARE_ROUNDED, tooltip="Share Public Link",
                          on_click=lambda e: page.run_task(on_share_report, e),
                          visible=len(state.charts) > 0),
        ],
    )

    return ft.View(route="/report", appbar=appbar, controls=[content], padding=0)
